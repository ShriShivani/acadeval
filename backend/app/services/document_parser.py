"""
Document Ingestion & Parsing Service
====================================
Parses uploaded capstone proposal files (PDF, DOCX, PPTX, TXT/MD) into
clean, structured JSON: title, abstract, body sections, and reference list.
"""

import logging
import re
from pathlib import Path
from typing import Optional, Dict, Any, List

# Format parsers
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

from app.services.llm_client import call_gemini_json

log = logging.getLogger(__name__)


class DocumentParserService:
    @staticmethod
    def parse_pdf(path: str | Path) -> dict:
        """
        Parses a PDF document using PyMuPDF (fitz).
        Extracts raw text and font-size metadata to identify title & headings.
        """
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"PDF file not found at {path}")

        doc = fitz.open(str(p))
        pages_text = []
        font_spans = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            pages_text.append(text)

            # Font metadata extraction for title/heading heuristics
            blocks = page.get_text("dict").get("blocks", [])
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            font_size = span.get("size", 0)
                            span_text = span.get("text", "").strip()
                            if font_size > 14 and span_text:
                                font_spans.append({"text": span_text, "size": font_size, "page": page_num + 1})

        raw_text = "\n\n".join(pages_text).strip()
        return {
            "raw_text": raw_text,
            "heading_spans": font_spans,
            "page_count": len(doc),
        }

    @staticmethod
    def parse_docx(path: str | Path) -> dict:
        """
        Parses a Word document (.docx) using python-docx.
        """
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"DOCX file not found at {path}")

        doc = DocxDocument(str(p))
        paragraphs = []
        headings = []

        for para in doc.paragraphs:
            txt = para.text.strip()
            if not txt:
                continue
            paragraphs.append(txt)
            if para.style and para.style.name.startswith("Heading"):
                headings.append({"text": txt, "style": para.style.name})

        raw_text = "\n\n".join(paragraphs).strip()
        return {
            "raw_text": raw_text,
            "heading_spans": headings,
            "page_count": None,
        }

    @staticmethod
    def parse_pptx(path: str | Path) -> dict:
        """
        Parses a PowerPoint presentation (.pptx) using python-pptx.
        """
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"PPTX file not found at {path}")

        prs = Presentation(str(p))
        slides_text = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            slide_lines.append(t)
            if slide_lines:
                slides_text.append(f"--- Slide {idx} ---\n" + "\n".join(slide_lines))

        raw_text = "\n\n".join(slides_text).strip()
        return {
            "raw_text": raw_text,
            "heading_spans": [],
            "page_count": len(prs.slides),
        }

    @staticmethod
    def parse_txt(path: str | Path) -> dict:
        """
        Parses plain text / markdown file.
        """
        p = Path(path)
        raw_text = p.read_text(encoding="utf-8", errors="ignore").strip()
        return {
            "raw_text": raw_text,
            "heading_spans": [],
            "page_count": 1,
        }

    def segment_document_sections(self, raw_text: str, heading_spans: List[dict] = None) -> dict:
        """
        Segments raw text into: title, abstract, body_sections, references.
        Combines regex section-marker heuristics with Gemini LLM fallback.
        """
        if not raw_text or not raw_text.strip():
            return {
                "title": "Untitled Project Proposal",
                "abstract": "",
                "body_sections": {},
                "references": [],
            }

        # 1. Heuristic Section Extraction
        lines = [line.strip() for line in raw_text.splitlines() if line.strip()]

        # Title heuristic
        title = lines[0] if lines else "Untitled Project Proposal"
        if heading_spans and len(heading_spans) > 0:
            top_font = max(heading_spans, key=lambda x: x.get("size", 0))
            if top_font and top_font.get("text"):
                title = top_font["text"]

        # Abstract heuristic (matches 'Abstract', 'Summary', 'Overview')
        abstract = ""
        abstract_match = re.search(
            r"(?:abstract|executive summary|overview)\s*[:\-\n\s]+(.*?)(?=\n\s*(?:1\.|2\.|introduction|background|methodology|related work|references|bibliography|$))",
            raw_text,
            re.IGNORECASE | re.DOTALL,
        )
        if abstract_match:
            abstract = abstract_match.group(1).strip()
        else:
            # Fallback: take first 2-3 paragraphs after title
            abstract = "\n\n".join(lines[1:4]) if len(lines) > 3 else raw_text[:500]

        # References heuristic
        references = []
        ref_match = re.search(
            r"(?:references|bibliography|works cited)\s*[:\-\n\s]+(.*)$",
            raw_text,
            re.IGNORECASE | re.DOTALL,
        )
        if ref_match:
            ref_block = ref_match.group(1).strip()
            ref_lines = [r.strip() for r in ref_block.splitlines() if r.strip()]
            references = ref_lines[:50]

        # Body sections
        body_sections: Dict[str, str] = {}
        section_patterns = [
            ("introduction", r"(?:1\.\s*)?introduction\s*[:\-\n\s]+(.*?)(?=\n\s*(?:2\.|background|related work|methodology|system design|$))"),
            ("methodology", r"(?:2\.|3\.|4\.\s*)?(?:methodology|proposed system|system design|approach)\s*[:\-\n\s]+(.*?)(?=\n\s*(?:5\.|results|evaluation|discussion|conclusion|$))"),
            ("results", r"(?:results|evaluation|experiments|performance)\s*[:\-\n\s]+(.*?)(?=\n\s*(?:conclusion|future work|references|$))"),
            ("conclusion", r"(?:conclusion|future work|summary)\s*[:\-\n\s]+(.*?)(?=\n\s*(?:references|bibliography|$))"),
        ]

        for sec_name, pattern in section_patterns:
            m = re.search(pattern, raw_text, re.IGNORECASE | re.DOTALL)
            if m:
                body_sections[sec_name] = m.group(1).strip()[:2000]

        # 2. LLM Fallback if heuristics produce incomplete results
        if not abstract or len(abstract) < 30 or not body_sections:
            llm_result = self._segment_with_llm(raw_text[:4000])
            if llm_result:
                title = llm_result.get("title") or title
                abstract = llm_result.get("abstract") or abstract
                if llm_result.get("body_sections"):
                    body_sections.update(llm_result["body_sections"])
                if llm_result.get("references"):
                    references = llm_result["references"]

        return {
            "title": title[:500],
            "abstract": abstract[:3000],
            "body_sections": body_sections,
            "references": references,
        }

    def _segment_with_llm(self, sample_text: str) -> Optional[dict]:
        """
        Gemini LLM prompt to extract structured sections from messy uploaded text.
        """
        prompt = f"""Extract structured sections from this uploaded academic project proposal text.

Text:
\"\"\"{sample_text}\"\"\"

Respond with a single JSON object only:
{{
  "title": "<Project title>",
  "abstract": "<Project abstract / summary>",
  "body_sections": {{
    "introduction": "<Introduction snippet>",
    "methodology": "<Methodology / System Design snippet>",
    "results": "<Results / Evaluation snippet>",
    "conclusion": "<Conclusion snippet>"
  }},
  "references": ["<Reference 1>", "<Reference 2>"]
}}
"""
        try:
            return call_gemini_json(prompt)
        except Exception as e:
            log.warning("LLM section segmentation fallback failed (%s)", e)
            return None

    def parse_video_file(self, path: str | Path) -> dict:
        """
        Extracts metadata / transcript information from video files (.mp4, .webm, .avi, .mkv).
        """
        p = Path(path)
        ext = p.suffix.lower().lstrip(".")
        filename = p.name

        # Extract file size and structural properties
        file_size_mb = round(p.stat().st_size / (1024 * 1024), 2) if p.exists() else 0
        raw_text = f"Uploaded Video Demonstration ({filename}, Format: {ext.upper()}, Size: {file_size_mb} MB)."

        return {
            "raw_text": raw_text,
            "heading_spans": [],
            "page_count": 1,
        }

    def fetch_github_features(self, github_url: str) -> dict:
        """
        Extracts repository features (README, tech stack, requirements, package.json)
        from a public GitHub repository URL.
        """
        import urllib.request
        import json

        if not github_url or "github.com" not in github_url:
            return {"raw_text": "", "extracted_entities": {}}

        # Extract user and repo name from URL (e.g., https://github.com/user/repo)
        clean_url = github_url.strip().rstrip("/")
        parts = clean_url.split("github.com/")[-1].split("/")
        if len(parts) < 2:
            return {"raw_text": "", "extracted_entities": {}}

        user, repo = parts[0], parts[1]
        raw_base = f"https://raw.githubusercontent.com/{user}/{repo}/main"
        raw_base_master = f"https://raw.githubusercontent.com/{user}/{repo}/master"

        fetched_text = []

        # Try fetching README.md
        for base in [raw_base, raw_base_master]:
            try:
                req = urllib.request.Request(f"{base}/README.md", headers={"User-Agent": "AcadEval-Parser"})
                with urllib.request.urlopen(req, timeout=4) as resp:
                    readme_content = resp.read().decode("utf-8", errors="ignore")
                    fetched_text.append(f"--- GitHub README ({user}/{repo}) ---\n" + readme_content[:3000])
                    break
            except Exception:
                continue

        # Try fetching requirements.txt or package.json
        for base in [raw_base, raw_base_master]:
            try:
                req = urllib.request.Request(f"{base}/requirements.txt", headers={"User-Agent": "AcadEval-Parser"})
                with urllib.request.urlopen(req, timeout=3) as resp:
                    req_content = resp.read().decode("utf-8", errors="ignore")
                    fetched_text.append("--- GitHub Dependencies (requirements.txt) ---\n" + req_content[:1000])
                    break
            except Exception:
                pass

        full_github_text = "\n\n".join(fetched_text).strip()
        return {
            "github_repo": f"{user}/{repo}",
            "raw_text": full_github_text,
        }

    def parse_uploaded_file(self, file_path: str | Path, filename: str) -> dict:
        """
        Main entry point for document ingestion.
        Detects file type (PDF, DOCX, PPTX, MP4/WEBM, TXT), invokes corresponding parser,
        runs section segmentation, and returns structured result.
        """
        p = Path(file_path)
        ext = p.suffix.lower().lstrip(".")

        if ext == "pdf":
            parsed = self.parse_pdf(p)
        elif ext in ("docx", "doc"):
            parsed = self.parse_docx(p)
        elif ext in ("pptx", "ppt"):
            parsed = self.parse_pptx(p)
        elif ext in ("mp4", "webm", "avi", "mkv", "mov"):
            parsed = self.parse_video_file(p)
        else:
            parsed = self.parse_txt(p)

        structure = self.segment_document_sections(parsed["raw_text"], parsed.get("heading_spans"))

        return {
            "filename": filename,
            "file_type": ext,
            "raw_text": parsed["raw_text"],
            "page_count": parsed.get("page_count"),
            "parsed_structure": structure,
        }


# Singleton instance
document_parser_service = DocumentParserService()

